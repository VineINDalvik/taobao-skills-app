#!/usr/bin/env python3
"""Generate competition PPTX for Tidal (潮汐) AI 全链路电商助手."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ─────────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x7C, 0x3A, 0xED)       # Purple (primary)
PRIMARY_LIGHT = RGBColor(0xED, 0xE9, 0xFE)  # Light purple bg
PRIMARY_DARK = RGBColor(0x4C, 0x1D, 0x95)   # Dark purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x18, 0x18, 0x1B)
GRAY = RGBColor(0x71, 0x71, 0x7A)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF5)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x25, 0x63, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(tf, text, size=16, color=BLACK, bold=False, space_before=Pt(4), align=PP_ALIGN.LEFT):
    """Add paragraph to existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.space_before = space_before
    p.alignment = align
    return p


def add_tag(slide, left, top, text, bg_color=PRIMARY_LIGHT, text_color=PRIMARY, size=11):
    """Add a small tag/badge."""
    w = Inches(max(1.2, len(text) * 0.18))
    shape = add_shape(slide, left, top, w, Inches(0.32), fill_color=bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY_DARK)

# Brand icon area
add_shape(slide, Inches(5.4), Inches(1.2), Inches(2.5), Inches(2.5), fill_color=PRIMARY)

add_text(slide, Inches(4.2), Inches(1.5), Inches(5), Inches(1.5),
         'Tidal', size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.2), Inches(2.7), Inches(5), Inches(0.6),
         '潮  汐', size=28, color=RGBColor(0xC4, 0xB5, 0xFD), bold=False, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.8),
         'AI 全链路电商助手', size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(4.7), Inches(8.3), Inches(0.6),
         '找爆款，让爆款持续爆', size=22, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(5.8), Inches(8.3), Inches(0.5),
         '面向淘宝中小商家（月销 < 100 万）的新品全链路落地工具', size=14, color=RGBColor(0xA7, 0x8B, 0xFA), align=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5),
         '淘宝开放平台 AI 技能大赛 · 2026', size=14, color=RGBColor(0xA7, 0x8B, 0xFA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Team & Business Scope
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.1 业务范围')
add_text(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.6),
         '团队介绍与核心能力', size=32, color=BLACK, bold=True)

# Team cards
team = [
    ('创始人 / 产品', 'DTC 跨境电商从业者\n转型淘宝中小商家服务', '电商运营全链路理解\n产品设计与用户洞察'),
    ('AI 工程', '视觉聚类模型实战经验\nGATv2 图神经网络', '选品预测、向量检索\n多模态 AI、因果推断'),
    ('前端工程', 'Next.js 全栈开发\nSaaS 产品 UI/UX', '产品原型快速落地\n现代化交互体验'),
]
for i, (role, bg_text, skill_text) in enumerate(team):
    x = Inches(0.8 + i * 4.0)
    card = add_shape(slide, x, Inches(1.8), Inches(3.6), Inches(2.2), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.3), Inches(1.95), Inches(3.0), Inches(0.4),
             role, size=18, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.4), Inches(3.0), Inches(0.7),
             bg_text, size=12, color=GRAY)
    add_text(slide, x + Inches(0.3), Inches(3.15), Inches(3.0), Inches(0.7),
             skill_text, size=12, color=BLACK)

# Core advantages
add_text(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
         '核心优势', size=22, color=BLACK, bold=True)

advantages = [
    ('真实 AI 选品模型', 'GATv2 + LightGBM 融合，S_F1=0.59（非 LLM 套壳）'),
    ('全链路串联', '7 个 Skill 数据自动流转，上游输出即下游输入'),
    ('闭环优化引擎', 'T+7 自动回测 + 指标追踪 + 建议迭代 = 持续优化飞轮'),
    ('运筹学定价', '运筹学 OR 模型，DID 因果推断 + 拉格朗日优化'),
]
for i, (title, desc) in enumerate(advantages):
    y = Inches(4.9 + i * 0.55)
    colors = [PRIMARY, BLUE, GREEN, AMBER]
    dot = add_shape(slide, Inches(1.0), y + Inches(0.06), Inches(0.15), Inches(0.15), fill_color=colors[i])
    add_text(slide, Inches(1.3), y, Inches(2.5), Inches(0.4),
             title, size=14, color=BLACK, bold=True)
    add_text(slide, Inches(3.8), y, Inches(8), Inches(0.4),
             desc, size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Market Pain Points
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.2 市场分析')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '1200 万中小商家的共同痛点', size=32, color=BLACK, bold=True)

# Big quote
quote_bg = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.0), fill_color=PRIMARY_LIGHT)
add_text(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(0.8),
         '"我刚进了一批货，接下来该怎么办？"', size=24, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# Competitor gaps
tools = [
    ('生意参谋', '给数据，不给决策', '商家看不懂也用不了', RED),
    ('万相台', '最低日预算 300+', '新品 0 销量无法冷启动', RED),
    ('店小蜜', '被动等问题', '不能主动帮商家卖货', AMBER),
    ('AIGC 工具', '只生成素材', '无选品/定价/运营整合', AMBER),
]
for i, (name, gap, detail, clr) in enumerate(tools):
    x = Inches(0.8 + i * 3.05)
    card = add_shape(slide, x, Inches(3.0), Inches(2.8), Inches(1.6), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), Inches(3.15), Inches(2.4), Inches(0.35),
             name, size=16, color=BLACK, bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(0.35),
             gap, size=13, color=clr, bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.9), Inches(2.4), Inches(0.5),
             detail, size=11, color=GRAY)

# Pain point stats table
add_text(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.4),
         '痛点量化', size=18, color=BLACK, bold=True)

pain_data = [
    ('不知道卖什么', '80%', '凭感觉/看排行', 'GATv2 AI 选品 + 趋势真假判断'),
    ('不会写标题/主图', '70%', '抄竞品', 'AI 生成 3 版标题 + 5 张主图方案'),
    ('不知道定多少钱', '65%', '看竞品定价', 'DID 弹性估算 + 拉格朗日最优价'),
    ('广告投了不赚钱', '60%', '盲投/关掉', '烧钱黑洞识别 + ROI 优化建议'),
    ('评价变差', '55%', '逐条看评价', 'P0/P1/P2 优先级 + 修复模板'),
]
# Header
headers = ['痛点', '影响面', '当前解法', 'Tidal 解法']
hx = [0.8, 4.0, 5.5, 8.0]
hw = [3.0, 1.3, 2.3, 4.5]
for j, h in enumerate(headers):
    add_text(slide, Inches(hx[j]), Inches(5.3), Inches(hw[j]), Inches(0.3),
             h, size=11, color=GRAY, bold=True)

for i, (pain, pct, curr, ours) in enumerate(pain_data):
    y = Inches(5.6 + i * 0.35)
    add_text(slide, Inches(0.8), y, Inches(3.0), Inches(0.3), pain, size=11, color=BLACK)
    add_text(slide, Inches(4.0), y, Inches(1.3), Inches(0.3), pct, size=11, color=RED, bold=True)
    add_text(slide, Inches(5.5), y, Inches(2.3), Inches(0.3), curr, size=11, color=GRAY)
    add_text(slide, Inches(8.0), y, Inches(4.5), Inches(0.3), ours, size=11, color=PRIMARY, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Market Size & Value
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.2 市场分析')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '业务空间与价值', size=32, color=BLACK, bold=True)

# Market size cards
mkt_cards = [
    ('1200 万+', '淘宝活跃商家', PRIMARY),
    ('90%+', '中小商家占比', BLUE),
    ('¥24 亿/年', '目标可触达市场', GREEN),
    ('10x+', '商家使用 ROI', AMBER),
]
for i, (num, label, clr) in enumerate(mkt_cards):
    x = Inches(0.8 + i * 3.05)
    card = add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(1.5), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), Inches(2.0), Inches(2.4), Inches(0.7),
             num, size=32, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.4),
             label, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Value sections
add_text(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(0.5),
         '对商家的价值', size=20, color=BLACK, bold=True)

merchant_vals = [
    '推广费从月均浪费 ¥2000 → 节省 ¥2000',
    '新品上架到起量周期缩短 50%+',
    '工具费 ¥199/月，ROI 超过 10 倍',
]
for i, v in enumerate(merchant_vals):
    add_text(slide, Inches(1.2), Inches(4.2 + i * 0.4), Inches(5.5), Inches(0.35),
             f'✓  {v}', size=14, color=BLACK)

add_text(slide, Inches(7.0), Inches(3.7), Inches(5.5), Inches(0.5),
         '对淘宝生态的价值', size=20, color=BLACK, bold=True)

eco_vals = [
    '降低中小商家经营门槛，提升平台活跃度',
    '丰富淘宝 AI 技能广场生态',
    '闭环数据飞轮——用户越多→模型越准→生态越活跃',
]
for i, v in enumerate(eco_vals):
    add_text(slide, Inches(7.4), Inches(4.2 + i * 0.4), Inches(5.5), Inches(0.35),
             f'✓  {v}', size=14, color=BLACK)

# Flywheel
fw_bg = add_shape(slide, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.5), fill_color=PRIMARY_LIGHT)
add_text(slide, Inches(2.8), Inches(5.6), Inches(7.7), Inches(0.4),
         '数据飞轮', size=16, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2.8), Inches(6.05), Inches(7.7), Inches(0.8),
         '用户越多 → T+7 回测数据越多 → 模型越准 → 推荐越好 → 用户越多',
         size=18, color=PRIMARY_DARK, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Product Overview — 7-Step Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '一个入口，七步闭环', size=32, color=BLACK, bold=True)
add_text(slide, Inches(0.8), Inches(1.45), Inches(10), Inches(0.4),
         '商家只需一个动作：描述选款方向 / 粘贴淘宝链接 / 上传灵感图', size=15, color=GRAY)

# Pipeline steps
steps = [
    ('🔍', 'AI 找款', '跟款', BLUE),
    ('🧪', '测款验证', '测款', PRIMARY),
    ('📝', '上架优化', '放量', GREEN),
    ('💰', '智能定价', '放量', GREEN),
    ('⭐', '评价诊断', '放量', GREEN),
    ('📊', '推广诊断', '放量', GREEN),
    ('🎯', '活动促销', '放量', GREEN),
]
for i, (icon, label, phase, clr) in enumerate(steps):
    x = Inches(0.6 + i * 1.75)
    # Icon circle
    circle = add_shape(slide, x + Inches(0.35), Inches(2.3), Inches(0.9), Inches(0.9), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.35), Inches(2.35), Inches(0.9), Inches(0.9),
             icon, size=28, align=PP_ALIGN.CENTER)
    # Label
    add_text(slide, x, Inches(3.3), Inches(1.6), Inches(0.35),
             label, size=13, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    # Phase tag
    phase_bg = RGBColor(0xDB, 0xEA, 0xFE) if phase == '跟款' else (PRIMARY_LIGHT if phase == '测款' else RGBColor(0xDC, 0xFC, 0xE7))
    phase_clr = BLUE if phase == '跟款' else (PRIMARY if phase == '测款' else GREEN)
    add_tag(slide, x + Inches(0.3), Inches(3.65), phase, bg_color=phase_bg, text_color=phase_clr, size=10)
    # Arrow
    if i < 6:
        add_text(slide, x + Inches(1.4), Inches(2.5), Inches(0.4), Inches(0.5),
                 '→', size=20, color=GRAY, align=PP_ALIGN.CENTER)

# Design principles
add_text(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.4),
         '核心设计原则', size=20, color=BLACK, bold=True)

principles = [
    ('输入极简', '最少：一张图 + 类目 + 成本价。商家最多填 3 个字段'),
    ('输出可执行', '每个输出都有 [复制] 按钮，一键贴进千牛后台。不是报告，是清单'),
    ('链路自动', '上游 Skill 输出自动成为下游输入，商家无感。从任意节点切入都可以'),
    ('闭环优化', 'T+7 自动回测效果，持续迭代建议。不是一次性报告，是优化飞轮'),
]
for i, (title, desc) in enumerate(principles):
    y = Inches(4.9 + i * 0.55)
    add_shape(slide, Inches(1.0), y + Inches(0.04), Inches(0.18), Inches(0.18), fill_color=PRIMARY)
    add_text(slide, Inches(1.4), y, Inches(2.0), Inches(0.4),
             title, size=14, color=PRIMARY, bold=True)
    add_text(slide, Inches(3.4), y, Inches(9), Inches(0.4),
             desc, size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Core Features — Skill 1 & 2
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '核心功能：AI 找款 + 测款验证', size=32, color=BLACK, bold=True)

# Skill 1 — left half
s1_bg = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.3), fill_color=LIGHT_GRAY)
add_text(slide, Inches(1.1), Inches(1.85), Inches(5.2), Inches(0.4),
         'Skill 1：AI 找款（爆款雷达）', size=18, color=PRIMARY, bold=True)
add_text(slide, Inches(1.1), Inches(2.3), Inches(5.2), Inches(0.35),
         '不知道卖什么 → 30 秒给出潜力爆款 + 设计基因解析', size=12, color=GRAY)

layers = [
    ('L1 风格聚类', '自研聚类引擎', 'FAISS 向量检索\n按视觉/情感相似度聚类', PRIMARY_LIGHT, PRIMARY),
    ('L2 趋势判断', '趋势过滤引擎', '区分真实上升趋势\nvs 短暂噪音', RGBColor(0xDB, 0xEA, 0xFE), BLUE),
    ('L3 爆款预测', '我们独有', 'GATv2 + LightGBM\nS/A/B 评分', RGBColor(0xFE, 0xF3, 0xC7), AMBER),
]
for i, (name, source, desc, bg, clr) in enumerate(layers):
    y = Inches(2.8 + i * 1.2)
    add_shape(slide, Inches(1.3), y, Inches(5.0), Inches(1.0), fill_color=bg)
    add_text(slide, Inches(1.5), y + Inches(0.05), Inches(2.2), Inches(0.35),
             name, size=13, color=clr, bold=True)
    add_text(slide, Inches(1.5), y + Inches(0.35), Inches(1.5), Inches(0.25),
             f'来自 {source}', size=9, color=GRAY)
    add_text(slide, Inches(3.8), y + Inches(0.1), Inches(2.3), Inches(0.8),
             desc, size=11, color=BLACK)

# Key output
add_text(slide, Inches(1.1), Inches(6.2), Inches(5.2), Inches(0.6),
         '输出：S/A/B 等级 + 趋势真实性 + 销量区间 + 设计基因 + 3 个延伸方向',
         size=11, color=PRIMARY)

# Skill 2 — right half
s2_bg = add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.3), fill_color=LIGHT_GRAY)
add_text(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.4),
         'Skill 2：测款验证', size=18, color=PRIMARY, bold=True)
add_text(slide, Inches(7.2), Inches(2.3), Inches(5.0), Inches(0.35),
         '从"凭感觉"到"数据驱动"的测款 SOP', size=12, color=GRAY)

test_items = [
    ('预测销量区间', '基于簇历史 P25/P50/P75'),
    ('8 维度数据录入', 'CTR / CVR / 加购率 / 收藏率等'),
    ('综合评分 0-100', '"建议放量" 或 "建议换款"'),
    ('测前 vs 测后对比', '模型自校准，越用越准'),
    ('止损规则', '花费上限 + 最低达标线'),
]
for i, (title, desc) in enumerate(test_items):
    y = Inches(2.9 + i * 0.85)
    add_shape(slide, Inches(7.2), y, Inches(0.35), Inches(0.35), fill_color=PRIMARY)
    add_text(slide, Inches(7.2), y, Inches(0.35), Inches(0.35),
             str(i+1), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.7), y, Inches(4.5), Inches(0.3),
             title, size=13, color=BLACK, bold=True)
    add_text(slide, Inches(7.7), y + Inches(0.35), Inches(4.5), Inches(0.3),
             desc, size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Core Features — Skill 3-7
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '核心功能：上架 · 定价 · 评价 · 推广 · 促销', size=32, color=BLACK, bold=True)

skills = [
    ('Skill 3：上架优化', '📝', [
        '3 版标题（搜索/种草/活动）+ CTR 预估',
        '5 张主图拍摄方案（构图/光线/场景/道具）',
        'AI 素材优化（去背/色彩/裁剪）',
        'SKU 推荐 + 关键词清单',
        '每个输出 [复制] → 千牛后台',
    ]),
    ('Skill 4：智能定价', '💰', [
        'DID 因果推断 → 真实价格弹性',
        '拉格朗日对偶 → GMV 最优价格',
        '20 档枚举表，最优行高亮',
        '三层风控：硬拦截/软警告/熔断',
        '阶段自动切换：新品→成长→大促→清尾',
    ]),
    ('Skill 5：评价诊断', '⭐', [
        'P0 立即处理（影响退货率）',
        'P1 本周处理（影响复购）',
        'P2 优化机会（好评词→标题）',
        '修复模板一键复制',
        '好评关键词反馈闭环',
    ]),
    ('Skill 6：推广诊断', '📊', [
        '烧钱黑洞识别（高花费低转化）',
        '加价机会（高 ROI 词加预算）',
        'ROI 优化预估',
        '基于利润上限计算最大可承受 PPC',
    ]),
    ('Skill 7：活动促销', '🎯', [
        '运筹学 OR 模型',
        '输入毛利目标 → 数学最优折扣',
        '活动节奏：预热/爆发/返场',
        '悲观场景预警 + 应急策略',
    ]),
]

cols = 5
for i, (title, icon, items) in enumerate(skills):
    x = Inches(0.5 + i * 2.5)
    w = Inches(2.3)
    card = add_shape(slide, x, Inches(1.7), w, Inches(5.3), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.15), Inches(1.85), w - Inches(0.3), Inches(0.6),
             f'{icon} {title}', size=12, color=PRIMARY, bold=True)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.15), Inches(2.5 + j * 0.55), w - Inches(0.3), Inches(0.5),
                 f'• {item}', size=10, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Product Screenshots
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '产品界面展示', size=32, color=BLACK, bold=True)
add_text(slide, Inches(0.8), Inches(1.45), Inches(10), Inches(0.4),
         '现代 SaaS 设计，对话式入口 + 多标签专业工作台', size=15, color=GRAY)

# 4 screenshot placeholders
screens = [
    ('首页：对话式入口', '自然语言描述 → AI 解析 → 确认选款\n支持链接解析 + JSON 导入\n快速开始样例 + 七步流程预览'),
    ('找款工作台', '簇卡片网格 + S/A/B 徽章\n站内外双视角趋势灵感\n三步引导：对齐簇→参考爆款→货源'),
    ('定价引擎', '算法步骤可视化\n20 档价格枚举表\n双轴折线图：利润 vs 日销'),
    ('商品管理看板', '健康分 0-100 + 进度条\n优化时间线（效果追踪）\n当前瓶颈 → 一键跳转 Skill'),
]
for i, (title, desc) in enumerate(screens):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(2.0 + (i // 2) * 2.7)
    # Screenshot placeholder
    placeholder = add_shape(slide, x, y, Inches(5.8), Inches(1.5),
                           fill_color=RGBColor(0xE4, 0xE4, 0xE7),
                           line_color=RGBColor(0xD4, 0xD4, 0xD8))
    add_text(slide, x + Inches(1.5), y + Inches(0.4), Inches(3), Inches(0.6),
             '[ 产品截图 ]', size=16, color=GRAY, align=PP_ALIGN.CENTER)
    # Label
    add_text(slide, x, y + Inches(1.55), Inches(5.8), Inches(0.3),
             title, size=13, color=BLACK, bold=True)
    add_text(slide, x, y + Inches(1.85), Inches(5.8), Inches(0.7),
             desc, size=10, color=GRAY)

add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
         '注：正式提交时请替换为实际产品截图（http://localhost:3000 各页面）', size=10, color=AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Evaluation Metrics
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '效果评估指标体系', size=32, color=BLACK, bold=True)

# Three layers
metric_layers = [
    ('模型层', PRIMARY, [
        ('GATv2 爆款预测准确率', 'S_F1 / Spearman', '≥ 0.59'),
        ('价格弹性估算误差', 'MAPE（T+7 回测）', '< 30%'),
        ('趋势真假判断', '30 天后验证', '> 75%'),
    ]),
    ('产品层', BLUE, [
        ('全链路完成率', '7 Skill 完成数', '> 60%'),
        ('操作时间节省', '对比人工', '> 70%'),
        ('建议采纳率', '执行/建议', '> 50%'),
    ]),
    ('商业层', GREEN, [
        ('商家 ROI 提升', '使用前后对比', '+30%+'),
        ('新品起量周期', '上架→日销>5', '缩短 50%'),
        ('推广费节省', '月均浪费减少', '> ¥2000'),
    ]),
]

for li, (layer_name, layer_color, metrics) in enumerate(metric_layers):
    x = Inches(0.5 + li * 4.2)
    add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(0.4), fill_color=layer_color)
    add_text(slide, x + Inches(0.2), Inches(1.72), Inches(3.5), Inches(0.35),
             layer_name, size=14, color=WHITE, bold=True)
    for mi, (name, method, target) in enumerate(metrics):
        y = Inches(2.2 + mi * 0.75)
        add_shape(slide, x, y, Inches(3.9), Inches(0.65), fill_color=LIGHT_GRAY)
        add_text(slide, x + Inches(0.15), y + Inches(0.03), Inches(2.5), Inches(0.3),
                 name, size=11, color=BLACK, bold=True)
        add_text(slide, x + Inches(0.15), y + Inches(0.33), Inches(2.0), Inches(0.25),
                 method, size=9, color=GRAY)
        add_text(slide, x + Inches(2.8), y + Inches(0.1), Inches(1.0), Inches(0.4),
                 target, size=14, color=layer_color, bold=True, align=PP_ALIGN.RIGHT)

# Feedback loop
add_text(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
         '闭环反馈机制', size=20, color=BLACK, bold=True)

fb_bg = add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.8), fill_color=PRIMARY_LIGHT)
fb_steps = [
    ('T+1 预警', '销量异常\n库存预警'),
    ('T+7 回测', '预测 vs 真实\n弹性校准'),
    ('T+30 校准', '更新训练集\n模型 fine-tune'),
    ('数据飞轮', '用户越多\n→ 模型越准'),
]
for i, (title, desc) in enumerate(fb_steps):
    x = Inches(1.2 + i * 3.0)
    add_text(slide, x, Inches(5.4), Inches(2.5), Inches(0.35),
             title, size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.8), Inches(2.5), Inches(0.8),
             desc, size=11, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + Inches(2.3), Inches(5.5), Inches(0.5), Inches(0.4),
                 '→', size=24, color=PRIMARY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Tech Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '技术架构与创新点', size=32, color=BLACK, bold=True)

# Tech stack — left
add_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
         '技术选型', size=18, color=BLACK, bold=True)

tech = [
    ('LLM', '通义千问 qwen-max', '淘宝生态、中文强'),
    ('图像理解', 'qwen-vl-max', '图片→属性识别'),
    ('向量检索', 'FAISS CPU', '无需部署，本地跑'),
    ('预测模型', 'GATv2 + LightGBM', '已有验证结果'),
    ('因果推断', 'DID 双重差分', '真实价格弹性'),
    ('优化求解', '拉格朗日 + 枚举', '数学最优定价'),
    ('前端', 'Next.js + Tailwind', '现代 SaaS 体验'),
]
for i, (layer, choice, reason) in enumerate(tech):
    y = Inches(2.1 + i * 0.5)
    add_text(slide, Inches(0.8), y, Inches(1.5), Inches(0.35),
             layer, size=11, color=GRAY, bold=True)
    add_text(slide, Inches(2.3), y, Inches(2.5), Inches(0.35),
             choice, size=11, color=BLACK, bold=True)
    add_text(slide, Inches(4.8), y, Inches(2.0), Inches(0.35),
             reason, size=10, color=GRAY)

# Innovations — right
add_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
         '五大核心创新', size=18, color=BLACK, bold=True)

innovations = [
    ('1', '真实 AI 选品模型', 'GATv2 图神经网络 + LightGBM\n非 LLM Prompt 套壳（S_F1=0.59）'),
    ('2', '三位一体选品', '风格聚类 + 趋势真假判断\n+ 爆款预测三层叠加'),
    ('3', '运筹学动态定价', '运筹学 OR 模型\nDID 因果推断 + 拉格朗日优化'),
    ('4', '评价→选品闭环', '差评标签→特征库\n选款自动过滤高风险属性'),
    ('5', '不滥用 LLM', '预测用小模型、优化用 OR\n因果推断用统计方法'),
]
for i, (num, title, desc) in enumerate(innovations):
    y = Inches(2.1 + i * 1.05)
    add_shape(slide, Inches(7.0), y, Inches(0.35), Inches(0.35), fill_color=PRIMARY)
    add_text(slide, Inches(7.0), y, Inches(0.35), Inches(0.35),
             num, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.5), y, Inches(5.0), Inches(0.3),
             title, size=13, color=BLACK, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.35), Inches(5.0), Inches(0.55),
             desc, size=10, color=GRAY)

# Model selection strategy bar
strat_bg = add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3), fill_color=LIGHT_GRAY)
add_text(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(0.3),
         '模型选择策略：用对工具，不滥用 LLM', size=14, color=BLACK, bold=True)

strategies = [
    ('结构化预测', 'GATv2 / LightGBM'),
    ('约束优化', 'OR 求解器'),
    ('向量检索', 'FAISS'),
    ('因果推断', 'DID'),
    ('语言生成', 'LLM'),
    ('图像理解', '多模态 LLM'),
]
for i, (problem, tool) in enumerate(strategies):
    x = Inches(1.0 + i * 2.0)
    add_text(slide, x, Inches(6.3), Inches(1.8), Inches(0.25),
             problem, size=10, color=GRAY)
    add_text(slide, x, Inches(6.55), Inches(1.8), Inches(0.25),
             f'→ {tool}', size=10, color=PRIMARY, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Inspiration Sources
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.3 创意方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '核心引擎矩阵', size=32, color=BLACK, bold=True)

sources = [
    ('自研聚类引擎', '大盘趋势分析', '降维到单商家：\n以图找款 + 设计基因解析', 'Skill 1'),
    ('趋势过滤引擎', '宏观趋势向导', '加入趋势真实性\n过滤层', 'Skill 1'),
    ('运筹学定价引擎', '大品牌大促运营', '降维到中小商家单品：\n阶段化动态定价', 'Skill 4+7'),
    ('自动对标引擎', '内部数据工具', '"当前款 vs 类目爆款"\n自动对标诊断', 'Skill 5'),
    ('延伸设计引擎', '资深买手经验', '爆款后自动生成\n3 个安全延伸方向', 'Skill 1'),
]
for i, (name, orig, adapt, skill) in enumerate(sources):
    x = Inches(0.5 + i * 2.5)
    card = add_shape(slide, x, Inches(1.7), Inches(2.3), Inches(3.5), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.15), Inches(1.85), Inches(2.0), Inches(0.7),
             name, size=12, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.15), Inches(2.6), Inches(2.0), Inches(0.25),
             f'原始：{orig}', size=9, color=GRAY)
    add_text(slide, x + Inches(0.15), Inches(2.95), Inches(2.0), Inches(0.3),
             '我们的改造：', size=9, color=BLACK, bold=True)
    add_text(slide, x + Inches(0.15), Inches(3.25), Inches(2.0), Inches(0.8),
             adapt, size=10, color=BLACK)
    add_tag(slide, x + Inches(0.15), Inches(4.5), skill, bg_color=PRIMARY_LIGHT, text_color=PRIMARY, size=10)

# Comparison bar
add_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
         '关键差异：我们不是单点工具', size=18, color=BLACK, bold=True)

# Coverage comparison
compare_bg = add_shape(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.0), fill_color=PRIMARY_LIGHT)
add_text(slide, Inches(1.0), Inches(6.15), Inches(11), Inches(0.35),
         '现有工具：选品 [ ] 测款 [ ] 上架 [部分] 定价 [ ] 评价 [数据] 推广 [付费] 促销 [ ]',
         size=12, color=GRAY)
add_text(slide, Inches(1.0), Inches(6.55), Inches(11), Inches(0.35),
         'Tidal：     选品 [AI] 测款 [AI] 上架 [AI]  定价 [AI] 评价 [AI]  推广 [AI]  促销 [AI]  ← 七环打通',
         size=12, color=PRIMARY, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Pricing Strategy
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.4 商业化方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '商业化定价策略', size=32, color=BLACK, bold=True)

tiers = [
    ('Free', '免费', [
        '每月 3 次全链路分析',
        '每次输出 1 个版本',
        '基础 7 Skill 功能',
    ], '用户获取 / 评审演示', LIGHT_GRAY, GRAY),
    ('Pro', '¥199/月', [
        '无限次全链路分析',
        '每次输出 3 版对比',
        '评价监控自动触发',
        'T+7 自动回测',
    ], '付费主力（单店商家）', PRIMARY_LIGHT, PRIMARY),
    ('Business', '¥599/月', [
        'Pro 全部功能',
        '多店铺（最多 5 家）',
        'API 接入（对接 ERP）',
        '专属顾问 1v1',
    ], '多店/团队商家', LIGHT_GRAY, GRAY),
]
for i, (name, price, features, target, bg, accent) in enumerate(tiers):
    x = Inches(0.8 + i * 4.1)
    w = Inches(3.8)
    is_pro = name == 'Pro'
    card = add_shape(slide, x, Inches(1.7), w, Inches(4.0),
                     fill_color=bg, line_color=PRIMARY if is_pro else None)
    if is_pro:
        add_tag(slide, x + Inches(2.3), Inches(1.8), '推荐', bg_color=PRIMARY, text_color=WHITE, size=10)
    add_text(slide, x + Inches(0.3), Inches(1.85), Inches(2.0), Inches(0.4),
             name, size=22, color=accent, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.3), w - Inches(0.6), Inches(0.5),
             price, size=28, color=BLACK, bold=True)
    for j, feat in enumerate(features):
        add_text(slide, x + Inches(0.3), Inches(2.95 + j * 0.4), w - Inches(0.6), Inches(0.35),
                 f'✓  {feat}', size=12, color=BLACK)
    add_text(slide, x + Inches(0.3), Inches(4.8 + (4 - len(features)) * 0.35), w - Inches(0.6), Inches(0.35),
             f'目标：{target}', size=10, color=GRAY)

# Justification
add_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.3),
         '定价合理性', size=16, color=BLACK, bold=True)
justifications = [
    '生意参谋专业版 ¥288/月（不给行动建议）→ 我们 ¥199 给完整行动清单',
    '帮商家月省推广费 ¥2000+ → 工具费 ¥199/月，ROI > 10 倍',
    '盈亏平衡：约 50 个 Pro 用户即可覆盖基础成本（LLM API + 服务器）',
    'LTV ¥1,888（年付）vs CAC ¥50-100 → LTV/CAC > 15',
]
for i, j in enumerate(justifications):
    add_text(slide, Inches(1.2), Inches(6.35 + i * 0.28), Inches(11), Inches(0.25),
             f'•  {j}', size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Scale Path
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.4 商业化方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '规模化路径', size=32, color=BLACK, bold=True)

# Three phases
phases = [
    ('第一阶段\n1-3 个月', '验证', [
        'AI 技能广场 + 商家群裂变',
        '小红书效果案例传播',
        '100 免费 → 20 付费用户',
        '核心：7 日留存 > 40%',
    ], BLUE),
    ('第二阶段\n3-12 个月', '增长', [
        '淘宝服务市场 + 培训机构',
        'SEM 投放',
        '1,000 付费用户',
        '月收入 ~¥20 万',
    ], PRIMARY),
    ('第三阶段\n1-2 年', '规模化', [
        '类目扩展：女装→男装→家居',
        '平台扩展：拼多多/京东',
        '数据飞轮加速',
        'ARR > 300 万',
    ], GREEN),
]
for i, (period, label, items, clr) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    w = Inches(3.8)
    add_shape(slide, x, Inches(1.7), w, Inches(0.35), fill_color=clr)
    add_text(slide, x + Inches(0.1), Inches(1.72), w - Inches(0.2), Inches(0.3),
             label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    card = add_shape(slide, x, Inches(2.1), w, Inches(2.8), fill_color=LIGHT_GRAY)
    add_text(slide, x + Inches(0.2), Inches(2.2), w - Inches(0.4), Inches(0.6),
             period, size=14, color=BLACK, bold=True)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.2), Inches(2.9 + j * 0.4), w - Inches(0.4), Inches(0.35),
                 f'•  {item}', size=12, color=BLACK)

# Revenue table
add_text(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
         '收入预测', size=18, color=BLACK, bold=True)

rev_bg = add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.5), fill_color=LIGHT_GRAY)
rev_headers = ['时间', '付费用户', 'ARPU', '月收入']
rev_data = [
    ('3 个月', '50', '¥199', '¥10,000'),
    ('6 个月', '300', '¥199', '¥60,000'),
    ('12 个月', '1,200', '¥220', '¥264,000'),
    ('24 个月', '5,000', '¥250', '¥1,250,000'),
]
for j, h in enumerate(rev_headers):
    x = Inches(1.0 + j * 2.8)
    add_text(slide, x, Inches(5.9), Inches(2.5), Inches(0.3),
             h, size=12, color=GRAY, bold=True)
for i, row in enumerate(rev_data):
    for j, val in enumerate(row):
        x = Inches(1.0 + j * 2.8)
        y = Inches(6.25 + i * 0.3)
        is_last_col = j == 3
        add_text(slide, x, y, Inches(2.5), Inches(0.25),
                 val, size=12, color=GREEN if is_last_col else BLACK,
                 bold=is_last_col)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Competitive Moats
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_tag(slide, Inches(0.8), Inches(0.5), '1.4 商业化方案')
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
         '竞争壁垒', size=32, color=BLACK, bold=True)

moats = [
    ('技术壁垒', 'GATv2 选品模型（已训练验证）\nDID 因果推断 + OR 定价\n需要大量电商数据和工程投入', '高', PRIMARY),
    ('数据壁垒', '评价→选品闭环 + T+7 回测\n数据飞轮：用户越多越准\n网络效应形成护城河', '高', GREEN),
    ('产品壁垒', '七环串联、链路自动流转\nvs 竞品单点工具\n全链路体验难以复制', '中高', BLUE),
    ('生态壁垒', '淘宝官方黑客松\n可获官方流量扶持\n先发优势 + 官方背书', '高', AMBER),
]
for i, (name, desc, level, clr) in enumerate(moats):
    x = Inches(0.5 + i * 3.15)
    w = Inches(2.9)
    card = add_shape(slide, x, Inches(1.7), w, Inches(3.5), fill_color=LIGHT_GRAY)
    add_shape(slide, x, Inches(1.7), w, Inches(0.5), fill_color=clr)
    add_text(slide, x + Inches(0.2), Inches(1.75), w - Inches(0.4), Inches(0.4),
             name, size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.35), w - Inches(0.4), Inches(2.0),
             desc, size=12, color=BLACK)
    add_tag(slide, x + Inches(0.2), Inches(4.5), f'可持续性：{level}', bg_color=clr, text_color=WHITE, size=10)

# Bottom summary
summary_bg = add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), fill_color=PRIMARY_LIGHT)
add_text(slide, Inches(1.2), Inches(5.7), Inches(11), Inches(0.4),
         '核心护城河：数据飞轮', size=18, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(6.15), Inches(11), Inches(0.6),
         '用户使用 → 产生真实交易数据 → T+7 回测 → 模型更准 → 推荐更好 → 更多用户使用\n这个飞轮一旦转起来，后来者需要同等量级的数据才能追平',
         size=14, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Closing
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY_DARK)

add_text(slide, Inches(2.5), Inches(1.5), Inches(8.3), Inches(1.5),
         'Tidal', size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(2.8), Inches(8.3), Inches(0.6),
         '潮  汐', size=28, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(0.8),
         '找爆款，让爆款持续爆', size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(5.0), Inches(8.3), Inches(0.5),
         '不只帮你找到爆款，更让它一直爆下去。', size=18, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

closing_points = [
    '一个入口，七步闭环',
    '真实 AI 模型，不是 LLM 套壳',
    '运筹学定价，数学最优',
    '数据飞轮，越用越准',
]
for i, p in enumerate(closing_points):
    add_text(slide, Inches(4.5), Inches(5.7 + i * 0.35), Inches(4.3), Inches(0.3),
             f'✓  {p}', size=14, color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.4),
         '提交入口：https://open.taobao.com/ai2026', size=12, color=RGBColor(0xA7, 0x8B, 0xFA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
out_dir = '/Users/vinexio/Desktop/dev-projects/taobao-skills-app/docs'
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'Tidal-AI-Competition-Deck.pptx')
prs.save(out_path)
print(f'Saved to {out_path}')
