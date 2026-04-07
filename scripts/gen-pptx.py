#!/usr/bin/env python3
"""Generate competition PPTX from the markdown content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ─────────────────────────────────────────────────
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)      # dark navy
ACCENT = RGBColor(0x3B, 0x82, 0xF6)       # blue-500
GREEN = RGBColor(0x16, 0xA3, 0x4A)        # green-600
RED = RGBColor(0xDC, 0x26, 0x26)          # red-600
AMBER = RGBColor(0xD9, 0x77, 0x06)        # amber-600
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=PRIMARY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=14, color=PRIMARY, bold=False, space_before=0, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=PRIMARY, line_spacing=1.2):
    """lines = list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        txt, bld = item[0], item[1] if len(item) > 1 else False
        clr = item[2] if len(item) > 2 else color
        sz = item[3] if len(item) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.space_before = Pt(4)
    return txBox

def section_tag(slide, left, top, text, bg_color=ACCENT):
    shape = add_rect(slide, left, top, Inches(2.2), Inches(0.35), bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_table_slide(slide, left, top, width, rows_data, col_widths=None):
    """rows_data = list of lists; first row is header."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.4 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = PRIMARY if r > 0 else WHITE
                p.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_BG
    return table_shape


# ════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "AI 全链路电商助手", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
         "找爆款，让爆款持续爆", 28, ACCENT, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
         "面向淘宝中小商家（月销 < 100 万）", 18, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
         "淘宝开放平台 AI 技能大赛 · 2026", 16, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 2: Team
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.1 业务范围")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "团队：DTC 实战背景 × AI 工程能力", 28, PRIMARY, True)

add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(12),
    [
        ["角色", "背景", "核心能力"],
        ["创始人 / 产品", "DTC 跨境电商从业者，正转型淘宝中小商家服务", "电商运营全链路理解、产品设计"],
        ["AI 工程", "视觉聚类模型（UMAP+HDBSCAN）、GATv2 图神经网络实战经验", "选品预测、向量检索、多模态 AI"],
        ["前端工程", "Next.js 全栈开发，SaaS 产品 UI/UX", "产品原型快速落地"],
    ])

add_text(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.5),
         "核心优势", 22, PRIMARY, True)
add_multiline(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(3), [
    ("1. 真实 AI 选品模型：GATv2 + LightGBM 融合模型（S_F1=0.59）——不是 LLM Prompt 套壳", False),
    ("2. 全链路串联：7 个 Skill 数据自动流转，上游输出即下游输入", False),
    ("3. 闭环优化引擎：T+7 自动回测 → 指标追踪 → 建议迭代——持续优化飞轮", False),
    ("4. 运筹学定价：DID 因果推断 + 拉格朗日对偶优化，专为中小商家打造", False),
    ("5. AI 实时生图：浏览器端 ONNX 前景分割 + GPT-image-1 图片生成/编辑", False),
], 13)


# ════════════════════════════════════════════════════════════
# SLIDE 3: Pain Points
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.2 市场分析")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "1200 万淘宝中小商家的共同痛点", 28, PRIMARY, True)
add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
         "“\u6211\u521a\u8fdb\u4e86\u4e00\u6279\u8d27\uff0c\u63a5\u4e0b\u6765\u8be5\u600e\u4e48\u529e\uff1f”", 20, ACCENT, True)

# Tool blind spots
add_multiline(slide, Inches(0.6), Inches(2.2), Inches(5.5), Inches(2), [
    ("现有工具的致命盲区：", True, RED),
    ("生意参谋 → 给数据，不给决策", False, GRAY),
    ("万相台 → 最低日预算 300+，新品无法冷启动", False, GRAY),
    ("AIGC 工具 → 只生成素材，无选品/定价/运营", False, GRAY),
], 13)

add_table_slide(slide, Inches(0.6), Inches(4.2), Inches(12),
    [
        ["痛点", "影响面", "当前解法", "我们的解法"],
        ["不知道卖什么", "80%", "凭感觉 / 看排行榜", "GATv2 AI 选品 + 趋势真假判断"],
        ["不会写标题/主图", "70%", "抄竞品", "AI 生成 3 版标题 + AI 改图流水线"],
        ["不知道定多少钱", "65%", "看竞品定价", "DID 弹性估算 + 拉格朗日最优价"],
        ["广告投了不赚钱", "60%", "盲投 / 关掉", "烧钱黑洞识别 + 关键词行动表"],
        ["评价变差不知道原因", "55%", "逐条看评价", "AI 买家画像 + 修复模板"],
        ["主图拍了但不好看", "50%", "反复试错", "AI 图片诊断 + 一键优化出图"],
    ])


# ════════════════════════════════════════════════════════════
# SLIDE 4: Market Size
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.2 市场分析")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "业务空间与价值", 28, PRIMARY, True)

# Market size cards
for i, (num, label) in enumerate([
    ("1200万+", "淘宝活跃商家"),
    ("90%+", "中小商家占比"),
    ("¥24亿/年", "目标可触达市场"),
    ("ROI > 10x", "商家投入产出比"),
]):
    left = Inches(0.6 + i * 3.1)
    card = add_rect(slide, left, Inches(1.7), Inches(2.8), Inches(1.4), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.6), num, 28, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.4), label, 13, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.6), Inches(3.5), Inches(6), Inches(0.5), "对商家", 20, PRIMARY, True)
add_multiline(slide, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.5), [
    ("• 推广费月均节省 2000 元", False),
    ("• 新品起量周期缩短 50%+", False),
    ("• 主图优化后 CTR 平均提升 35%+", False),
    ("• 工具费 ¥199/月，ROI > 10 倍", False, GREEN),
], 14)

add_text(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(0.5), "对淘宝生态", 20, PRIMARY, True)
add_multiline(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(2.5), [
    ("• 降低中小商家经营门槛", False),
    ("• 丰富淘宝 AI 技能广场生态", False),
    ("• 闭环数据飞轮加速平台活跃度", False),
], 14)


# ════════════════════════════════════════════════════════════
# SLIDE 5: Product Overview (7-step)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.3 创意方案")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "一个入口，七步闭环", 28, PRIMARY, True)
add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
         "商家只需一个动作：上传商品图片 / 输入链接 / 描述选款方向 → AI 自动完成全部 7 步", 16, GRAY)

# 7 skill boxes in a flow
skills = [
    ("1", "AI 找款", "GATv2 选品\n趋势判断"),
    ("1.5", "测款验证", "预测区间\n综合评分"),
    ("2", "上架优化", "标题/主图\nAI 改图"),
    ("3", "智能定价", "DID+拉格朗日\n20档枚举"),
    ("4", "评价诊断", "买家画像\n修复模板"),
    ("5", "推广诊断", "烧钱识别\n关键词优化"),
    ("6", "活动促销", "OR 最优折扣\n排期规划"),
]
for i, (num, name, desc) in enumerate(skills):
    left = Inches(0.3 + i * 1.8)
    card = add_rect(slide, left, Inches(2.2), Inches(1.6), Inches(2.2), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.1), Inches(2.3), Inches(1.4), Inches(0.35),
             f"Skill {num}", 11, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), Inches(2.65), Inches(1.4), Inches(0.4),
             name, 15, PRIMARY, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), Inches(3.1), Inches(1.4), Inches(0.9),
             desc, 10, GRAY, False, PP_ALIGN.CENTER)
    # Arrow between cards
    if i < len(skills) - 1:
        add_text(slide, left + Inches(1.55), Inches(3.0), Inches(0.3), Inches(0.3),
                 "→", 16, ACCENT, True, PP_ALIGN.CENTER)

# Bottom: management board
add_rect(slide, Inches(3), Inches(4.8), Inches(7), Inches(1), LIGHT_BG, GREEN)
add_text(slide, Inches(3.2), Inches(4.9), Inches(6.6), Inches(0.4),
         "商品管理看板", 16, GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3.2), Inches(5.3), Inches(6.6), Inches(0.4),
         "健康分 + 优化时间线 + 当前瓶颈 → 一键跳转 Skill + 会话恢复", 12, GRAY, False, PP_ALIGN.CENTER)

# Design principles
add_multiline(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(1.2), [
    ("核心设计原则：输入越少越好 → 输出越可执行越好 → 7 个 Skill 对商家透明，感知的是'一个助手'", False, GRAY, 12),
], 12)


# ════════════════════════════════════════════════════════════
# SLIDE 6: Skill 1 + 1.5 Details
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "Skill 1: AI 找款")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "三层处理架构 + 延伸设计引擎", 28, PRIMARY, True)

# 3-layer architecture
add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(6),
    [
        ["层级", "引擎", "功能"],
        ["L1 风格聚类", "自研聚类引擎", "FAISS + OpenCLIP 视觉嵌入聚类"],
        ["L2 趋势判断", "趋势过滤引擎", "区分真实趋势 vs 短暂噪音"],
        ["L3 爆款预测", "我们独有", "GATv2 + LightGBM → S/A/B"],
    ])

# Extension design engine
add_text(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
         "延伸设计实验室", 18, PRIMARY, True)
add_multiline(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(2.5), [
    ("浏览器端 ONNX 前景分割 → 实时提取服装主体", False),
    ("Canvas 蒙版叠加：5 种预设色 + 4 种图案", False),
    ("两种改款模式：", True),
    ("  ① 提示词改款：自然语言 → AI 生图", False),
    ("  ② 设计点库：元素/面料/廓形 → 精准 prompt", False),
    ("真实调用 GPT-image-1 生成图片", False, GREEN),
    ("基于簇语义标签自动生成编辑 prompt", False, GREEN),
], 12)

# Testing
section_tag(slide, Inches(0.6), Inches(4.0), "Skill 1.5: 测款验证")
add_multiline(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), [
    ("预测销量区间 P25/P50/P75（基于簇历史数据）+ S/A/B 等级", False, PRIMARY, 13),
    ("8 维度测款数据录入（CTR/CVR/加购率/收藏率等）", False, PRIMARY, 13),
    ("小模型回调 + LLM 叙事融合 → 综合评分（0-100）→ '建议放量' / '建议换款'", True, ACCENT, 13),
    ("维度对比条 + 放量预测（月营收/月利润/回本天数）", False, PRIMARY, 13),
    ("测款数据匿名提交到簇级统计池 → 反馈模型 → 数据飞轮", False, GREEN, 13),
], 13)


# ════════════════════════════════════════════════════════════
# SLIDE 7: Skill 2 (Listing) Details
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "Skill 2: 上架优化")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "一键生成上架包 + AI 四步改图 + 图片诊断", 26, PRIMARY, True)

# 4 tabs
tabs = [
    ("待办清单", "进度仪表盘\nAI 建议优先操作\n跨标签跳转"),
    ("标题·详情·SKU", "3版标题 + CTR预估\n详情分屏文案\nSKU + 直通车词"),
    ("AI改图·拍摄简报", "四步改图流水线\n5张拍摄简报\n状态追踪"),
    ("图片诊断", "综合评分 0-100\n爆款套路对标\n改图优先级"),
]
for i, (name, desc) in enumerate(tabs):
    left = Inches(0.6 + i * 3.1)
    card = add_rect(slide, left, Inches(1.7), Inches(2.8), Inches(1.5), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.15), Inches(1.8), Inches(2.5), Inches(0.35), name, 14, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.15), Inches(2.5), Inches(1.0), desc, 11, GRAY, False, PP_ALIGN.CENTER)

# AI 4-step pipeline
add_text(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
         "AI 四步改图流水线", 18, PRIMARY, True)
steps = ["① 找问题：上传图 → AI 检测问题", "② 看原因：每个问题对销量的影响",
         "③ 选方案：勾选改进项（至少 3 项）", "④ 出新图：AI 生成 3 套优化方案"]
for i, step in enumerate(steps):
    left = Inches(0.6 + i * 3.1)
    add_rect(slide, left, Inches(4.0), Inches(2.8), Inches(0.7), LIGHT_BG, GREEN)
    add_text(slide, left + Inches(0.1), Inches(4.05), Inches(2.6), Inches(0.6), step, 11, PRIMARY, False, PP_ALIGN.CENTER)

# Scoring
add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
         "图片诊断", 18, PRIMARY, True)
add_multiline(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(2), [
    ("综合评分环形图（构图/色彩/文案/版式）→ 3 种爆款拍图套路 → 你的图 vs 爆款对比表 → 改图优先级排序", False, PRIMARY, 13),
    ("拍摄简报：5 张 × P0/P1/P2 优先级，每张含场景/光线/构图/道具/模特要求，一键复制拍摄单", False, GRAY, 12),
], 13)


# ════════════════════════════════════════════════════════════
# SLIDE 8: Skill 3 (Pricing) Details
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "Skill 3: 智能定价")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "运筹学动态定价引擎", 28, PRIMARY, True)

# 4-step algorithm
algo_steps = [
    ("DID 因果推断", "剥离季节性\n得出真实弹性"),
    ("LightGBM 冷启动", "新品无数据时\n弹性预估"),
    ("拉格朗日优化", "毛利约束下\n求 GMV 最大价"),
    ("三层风控护栏", "负毛利拦截\n深降价警告\n日销熔断"),
]
for i, (name, desc) in enumerate(algo_steps):
    left = Inches(0.6 + i * 3.1)
    card = add_rect(slide, left, Inches(1.7), Inches(2.8), Inches(1.8), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.15), Inches(1.85), Inches(2.5), Inches(0.35), name, 15, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.3), Inches(2.5), Inches(1.0), desc, 12, GRAY, False, PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, left + Inches(2.75), Inches(2.3), Inches(0.4), Inches(0.4), "→", 20, ACCENT, True, PP_ALIGN.CENTER)

add_multiline(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(3), [
    ("20 个价格档位枚举表：日销/GMV/月利润对比，最优行高亮", False),
    ("阶段化推荐：新品期 → 成长期 → 大促期 → 清尾期 自动切换", False),
    ("Recharts 双轴折线图：月利润 vs 日销量可视化", False),
    ("数据流入：测款成本价 → 毛利约束", False, GREEN),
    ("数据流出：价格弹性 → 推广 PPC 出价上限 / 日常价 → 促销折扣锚定", False, GREEN),
], 14)


# ════════════════════════════════════════════════════════════
# SLIDE 9: Skill 4-6 + Management
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "Skill 4-6 + 管理看板")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "评价诊断 · 推广诊断 · 活动促销 · 商品管理", 26, PRIMARY, True)

# Skill 4
add_rect(slide, Inches(0.6), Inches(1.7), Inches(3.8), Inches(2.6), LIGHT_BG, ACCENT)
add_text(slide, Inches(0.8), Inches(1.8), Inches(3.4), Inches(0.35), "Skill 4: 评价诊断", 15, ACCENT, True)
add_multiline(slide, Inches(0.8), Inches(2.2), Inches(3.4), Inches(2), [
    ("AI 买家画像深潜", True, PRIMARY, 12),
    ("多维度画像 + 拟真评价", False, GRAY, 11),
    ("产品匹配度评分（0-10）", False, GRAY, 11),
    ("P0/P1/P2 修复行动 + 模板", False, GRAY, 11),
    ("好评词 → 回传标题优化", False, GREEN, 11),
], 11)

# Skill 5
add_rect(slide, Inches(4.7), Inches(1.7), Inches(3.8), Inches(2.6), LIGHT_BG, ACCENT)
add_text(slide, Inches(4.9), Inches(1.8), Inches(3.4), Inches(0.35), "Skill 5: 推广诊断", 15, ACCENT, True)
add_multiline(slide, Inches(4.9), Inches(2.2), Inches(3.4), Inches(2), [
    ("烧钱黑洞识别", True, PRIMARY, 12),
    ("关键词行动表：暂停/加价/新增", False, GRAY, 11),
    ("最大 PPC = 单价×CVR×(1/ROI)", False, GRAY, 11),
    ("投放高峰时段建议", False, GRAY, 11),
    ("优化后预估 ROI 提升", False, GREEN, 11),
], 11)

# Skill 6
add_rect(slide, Inches(8.8), Inches(1.7), Inches(3.8), Inches(2.6), LIGHT_BG, ACCENT)
add_text(slide, Inches(9.0), Inches(1.8), Inches(3.4), Inches(0.35), "Skill 6: 活动促销", 15, ACCENT, True)
add_multiline(slide, Inches(9.0), Inches(2.2), Inches(3.4), Inches(2), [
    ("运筹学 OR 最优折扣", True, PRIMARY, 12),
    ("预热/爆发/返场执行排期", False, GRAY, 11),
    ("Monte Carlo 三情景利润预测", False, GRAY, 11),
    ("保存到商品库 + JSON 导出", False, GRAY, 11),
], 11)

# Management board
add_rect(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2.3), LIGHT_BG, GREEN)
add_text(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.35), "商品管理看板", 18, GREEN, True)
add_multiline(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), [
    ("商品卡片：健康分环形图 + Skills 完成进度条 + 关键指标", False, PRIMARY, 13),
    ("优化时间线：每次操作（改标题/调价/换主图）效果追踪 → 标记有效/待观察/已回滚", False, PRIMARY, 13),
    ("当前瓶颈提示 → 一键跳转 Skill + 会话恢复（从上次完成的 Skill 继续）", False, PRIMARY, 13),
], 13)


# ════════════════════════════════════════════════════════════
# SLIDE 10: Evaluation Metrics
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.3 效果评估")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "三层评估指标体系 + 闭环反馈", 28, PRIMARY, True)

add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(12),
    [
        ["层级", "指标", "衡量方式", "目标值"],
        ["模型层", "GATv2 爆款预测准确率", "S_F1 / Spearman", "S_F1 ≥ 0.59"],
        ["模型层", "价格弹性估算误差", "MAPE（T+7 回测）", "< 30%"],
        ["模型层", "图片诊断评分一致性", "与人工评审对比", "Kappa > 0.6"],
        ["产品层", "全链路完成率", "完成数/总数", "> 60%"],
        ["产品层", "AI 改图采纳率", "替换上线/生成总数", "> 30%"],
        ["商业层", "商家 ROI 提升", "使用前后对比", "+30%+"],
        ["商业层", "主图优化后 CTR 提升", "替换前后 A/B", "+35%+"],
        ["商业层", "推广费节省", "月均浪费减少", "> ¥2000/月"],
    ])

add_text(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.4),
         "数据飞轮：用户越多 → 回测数据越多 → 模型越准 → 推荐越好 → 用户越多", 16, GREEN, True)
add_text(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
         "闭环：T+1 预警 → T+7 回测 → T+30 全量校准（测款数据匿名回馈簇级统计池）", 14, GRAY)


# ════════════════════════════════════════════════════════════
# SLIDE 11: Tech Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.3 技术架构")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "技术选型与六大核心创新", 28, PRIMARY, True)

add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(5.8),
    [
        ["层", "选型", "理由"],
        ["LLM", "通义千问 qwen-max", "淘宝生态、中文强"],
        ["图片生成", "GPT-image-1 / mini", "蒙版 + 参考图编辑"],
        ["图片分割", "浏览器端 ONNX", "零服务端成本"],
        ["向量检索", "FAISS + OpenCLIP", "高效无推理"],
        ["预测模型", "GATv2 + LightGBM", "已有验证"],
        ["因果推断", "DID 双重差分", "真实弹性"],
        ["优化求解", "拉格朗日对偶", "数学最优"],
        ["前端", "Next.js 16 + Zustand", "现代 SaaS"],
    ])

add_text(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
         "六大核心创新", 18, PRIMARY, True)
innovations = [
    ("1", "真实 AI 选品", "GATv2+LightGBM\n非 LLM 套壳"),
    ("2", "三位一体聚类", "风格+趋势+预测\n三层叠加"),
    ("3", "运筹学定价", "OR 模型\n阶段自动切换"),
    ("4", "AI 图片闭环", "ONNX分割→预览\n→GPT生成→诊断"),
    ("5", "买家画像深潜", "AI画像+拟真评价\n+匹配度评分"),
    ("6", "不滥用 LLM", "小模型/OR/统计\n各取所长"),
]
for i, (num, name, desc) in enumerate(innovations):
    row = i // 3
    col = i % 3
    left = Inches(7 + col * 1.95)
    top = Inches(2.2 + row * 2.3)
    card = add_rect(slide, left, top, Inches(1.8), Inches(2.0), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.1), top + Inches(0.1), Inches(1.6), Inches(0.3),
             f"#{num}", 11, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.4), Inches(1.6), Inches(0.35),
             name, 12, PRIMARY, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.8), Inches(1.6), Inches(1.0),
             desc, 10, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 12: Cross-Skill Data Flow
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.3 核心引擎")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "六大核心引擎 + 跨 Skill 数据流转", 28, PRIMARY, True)

add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(12),
    [
        ["核心引擎", "原始定位", "我们的改造", "落地"],
        ["自研聚类引擎", "大盘趋势分析", "降维到单商家：以图找款 + 设计基因", "Skill 1"],
        ["趋势过滤引擎", "宏观趋势向导", "加入趋势真实性过滤层", "Skill 1"],
        ["运筹学定价引擎", "大品牌大促", "降维到单品：阶段化动态定价", "Skill 3+6"],
        ["延伸设计引擎", "资深买手经验", "自动生成延伸方向 + AI 实时生图", "Skill 1"],
        ["AI 素材优化引擎", "专业设计师审图", "商家自助：四步改图 + 诊断 + 对标", "Skill 2"],
    ])

add_text(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
         "跨 Skill 数据流转（非单点工具，七环打通）", 16, GREEN, True)
add_multiline(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(2.2), [
    ("Skill 1 → 设计基因/簇语义 → Skill 2 标题含款式关键词", False),
    ("Skill 1 → S/A/B 等级 → Skill 1.5 预测区间基准", False),
    ("Skill 1.5 → 成本价 → Skill 3 毛利约束", False),
    ("Skill 3 → 价格弹性 → Skill 5 PPC 出价上限  |  日常价 → Skill 6 折扣锚定", False),
    ("Skill 4 → 好评关键词 → Skill 2 标题优化  |  差评标签 → Skill 1 过滤高风险", False, GREEN),
    ("Skill 5 → 真实 ROI → Skill 3 触发重新定价建议", False, GREEN),
], 12)


# ════════════════════════════════════════════════════════════
# SLIDE 13: Pricing Strategy
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.4 商业化方案")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "三级定价策略", 28, PRIMARY, True)

add_table_slide(slide, Inches(0.6), Inches(1.7), Inches(12),
    [
        ["套餐", "价格", "核心权益", "目标用户"],
        ["Free", "免费", "每月 3 次全链路分析", "新用户获取"],
        ["Pro", "¥199/月 | ¥1,888/年", "无限分析 + 评价监控 + AI 改图无限次", "单店商家（主力）"],
        ["Business", "¥599/月", "Pro + 多店铺(5家) + API + 批量 AI 生图", "多店/团队"],
    ])

add_text(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
         "定价合理性", 18, PRIMARY, True)
add_multiline(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(2.5), [
    ("生意参谋专业版 ¥288/月（不给行动建议）", False, GRAY),
    ("万相台月均消耗 ¥5,000-20,000", False, GRAY),
    ("美工外包改图 ¥50-200/张", False, GRAY),
    ("我们帮商家省的 > 工具费（ROI > 10x）", True, GREEN),
    ("盈亏平衡：约 50 个 Pro 用户", False),
], 14)

# Revenue projection
add_table_slide(slide, Inches(7), Inches(3.5), Inches(5.5),
    [
        ["时间", "付费用户", "ARPU", "月收入"],
        ["3 个月", "50", "¥199", "¥10,000"],
        ["6 个月", "300", "¥199", "¥60,000"],
        ["12 个月", "1,200", "¥220", "¥264,000"],
        ["24 个月", "5,000", "¥250", "¥1,250,000"],
    ])
add_text(slide, Inches(7), Inches(5.9), Inches(5.5), Inches(0.4),
         "LTV ¥1,888 | CAC ¥50-100 | LTV/CAC > 15", 13, GREEN, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 14: GTM + Barriers
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_tag(slide, Inches(0.6), Inches(0.4), "1.4 商业化方案")
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6),
         "规模化路径 + 竞争壁垒", 28, PRIMARY, True)

# GTM 3 phases
phases = [
    ("第一阶段 1-3月", "验证", "AI 技能广场 + KOC 裂变\n100 免费 → 20 付费\n7 日留存 > 40%"),
    ("第二阶段 3-12月", "增长", "淘宝服务市场 + 培训合作\n1,000 付费（月收¥20万）\nNPS > 50"),
    ("第三阶段 1-2年", "规模化", "类目扩展 + 平台扩展\n5,000 付费（月收¥125万）\n数据飞轮加速"),
]
for i, (time, label, desc) in enumerate(phases):
    left = Inches(0.6 + i * 4.1)
    card = add_rect(slide, left, Inches(1.7), Inches(3.8), Inches(2), LIGHT_BG, ACCENT)
    add_text(slide, left + Inches(0.15), Inches(1.8), Inches(3.5), Inches(0.3), time, 12, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.1), Inches(3.5), Inches(0.3), label, 18, PRIMARY, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), Inches(2.5), Inches(3.5), Inches(1.0), desc, 11, GRAY, False, PP_ALIGN.CENTER)

# Barriers
add_text(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
         "四层竞争壁垒", 18, PRIMARY, True)
add_table_slide(slide, Inches(0.6), Inches(4.5), Inches(12),
    [
        ["壁垒", "内容", "可持续性"],
        ["技术壁垒", "GATv2 选品 + DID 因果推断 + OR 定价 + ONNX 分割 + GPT-image-1 改图", "高"],
        ["数据壁垒", "评价→选品闭环 + T+7 回测飞轮 + 测款回馈池", "高"],
        ["产品壁垒", "七环串联 + AI 改图诊断全链路（vs 竞品单点工具）", "中"],
        ["生态壁垒", "淘宝官方黑客松，可获官方流量扶持", "高"],
    ])


# ════════════════════════════════════════════════════════════
# SLIDE 15: Closing
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
         "AI 全链路电商助手", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.6),
         "不只帮你找到爆款，更让爆款持续爆下去", 24, ACCENT, True, PP_ALIGN.CENTER)

bullets = [
    "一个入口，七步闭环",
    "真实 AI 模型，不是 LLM 套壳",
    "运筹学定价，数学最优",
    "AI 改图 + 图片诊断，商家自助出好图",
    "数据飞轮，越用越准",
]
for i, b in enumerate(bullets):
    add_text(slide, Inches(4), Inches(3.5 + i * 0.45), Inches(5), Inches(0.4),
             f"•  {b}", 16, WHITE, False, PP_ALIGN.LEFT)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
         "让 1200 万中小商家，都能用上专业运营团队级的 AI 能力", 18, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
         "淘宝开放平台 AI 技能大赛 · 2026", 14, GRAY, False, PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────
output_path = "/Users/vinexio/Desktop/dev-projects/taobao-skills-app/docs/competition-ppt.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
